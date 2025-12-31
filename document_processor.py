# document_processor.py

import fitz  # PyMuPDF
from pptx import Presentation
import re
from typing import List, Dict
import tiktoken

class DocumentProcessor:
    """
    Procesa documentos médicos (PDF y PPT) y los convierte en chunks
    listos para indexar en el sistema RAG
    """
    
    def __init__(self):
        # Tokenizer para contar tokens (compatible con GPT/Claude/Gemini)
        self.encoder = tiktoken.get_encoding("cl100k_base")
    
    def count_tokens(self, text: str) -> int:
        """Cuenta tokens en un texto"""
        return len(self.encoder.encode(text))
    
    def extract_from_pdf(self, pdf_path: str) -> Dict:
        """
        Extrae texto de un PDF manteniendo estructura
        
        Args:
            pdf_path: ruta al archivo PDF
        
        Returns:
            dict con 'metadata' y 'chunks'
        """
        print(f"📄 Abriendo PDF: {pdf_path}")
        
        try:
            doc = fitz.open(pdf_path)
        except Exception as e:
            print(f"❌ Error al abrir PDF: {e}")
            return {"metadata": {}, "chunks": []}
        
        chunks = []
        
        # Extraer metadatos
        metadata = {
            "title": self._extract_title(doc),
            "pages": len(doc),
            "type": "guideline"
        }
        
        print(f"📖 Título detectado: {metadata['title']}")
        print(f"📄 Total páginas: {len(doc)}")
        
        # Procesar cada página
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")
            
            # Limpiar texto
            text = self._clean_text(text)
            
            if not text.strip():
                continue
            
            # Dividir en chunks si es necesario
            sections = self._split_by_sections(text, page_num + 1)
            chunks.extend(sections)
            
            if (page_num + 1) % 10 == 0:
                print(f"  ✓ Procesadas {page_num + 1} páginas...")
        
        print(f"✅ PDF procesado: {len(chunks)} chunks extraídos")
        
        doc.close()
        
        return {
            "metadata": metadata,
            "chunks": chunks
        }
    
    def extract_from_ppt(self, ppt_path: str) -> Dict:
        """
        Extrae texto de un PowerPoint
        
        Args:
            ppt_path: ruta al archivo PPTX
        
        Returns:
            dict con 'metadata' y 'chunks'
        """
        print(f"📊 Abriendo PPT: {ppt_path}")
        
        try:
            prs = Presentation(ppt_path)
        except Exception as e:
            print(f"❌ Error al abrir PPT: {e}")
            return {"metadata": {}, "chunks": []}
        
        chunks = []
        
        # Metadatos
        metadata = {
            "title": prs.core_properties.title or "Presentación sin título",
            "pages": len(prs.slides),
            "type": "presentation"
        }
        
        print(f"📖 Título: {metadata['title']}")
        print(f"📊 Total slides: {len(prs.slides)}")
        
        # Procesar cada slide
        for slide_num, slide in enumerate(prs.slides, 1):
            # Extraer título del slide
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            
            # Extraer todo el contenido
            content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content.append(shape.text)
            
            # Combinar título + contenido
            full_text = f"{title}\n\n" + "\n".join(content)
            full_text = self._clean_text(full_text)
            
            if not full_text.strip():
                continue
            
            chunks.append({
                "text": full_text,
                "page": slide_num,
                "section": title or f"Slide {slide_num}",
                "type": "slide",
                "tokens": self.count_tokens(full_text)
            })
            
            if slide_num % 10 == 0:
                print(f"  ✓ Procesados {slide_num} slides...")
        
        print(f"✅ PPT procesado: {len(chunks)} slides extraídos")
        
        return {
            "metadata": metadata,
            "chunks": chunks
        }
    
    def _clean_text(self, text: str) -> str:
        """
        Limpia texto manteniendo estructura médica importante
        """
        # Normalizar espacios múltiples
        text = re.sub(r' +', ' ', text)
        
        # Normalizar saltos de línea múltiples
        text = re.sub(r'\n\n+', '\n\n', text)
        
        # Quitar espacios al inicio/final
        text = text.strip()
        
        return text
    
    def _split_by_sections(self, text: str, page_num: int) -> List[Dict]:
        """
        Divide texto en chunks inteligentes
        Respeta la estructura médica (no corta en medio de una lista, tabla, etc.)
        """
        # Si es corto, devolver completo
        if self.count_tokens(text) < 600:
            return [{
                "text": text,
                "page": page_num,
                "section": f"Página {page_num}",
                "tokens": self.count_tokens(text)
            }]
        
        # Dividir por párrafos
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        chunks = []
        current_chunk = ""
        current_section = f"Página {page_num}"
        
        # Intentar detectar título de sección (línea corta seguida de contenido)
        if paragraphs and len(paragraphs[0]) < 100:
            current_section = paragraphs[0][:100]
        
        for para in paragraphs:
            # Intentar agregar al chunk actual
            test_chunk = current_chunk + "\n\n" + para if current_chunk else para
            
            if self.count_tokens(test_chunk) < 600:
                current_chunk = test_chunk
            else:
                # Guardar chunk actual
                if current_chunk:
                    chunks.append({
                        "text": current_chunk.strip(),
                        "page": page_num,
                        "section": current_section,
                        "tokens": self.count_tokens(current_chunk)
                    })
                
                # Empezar nuevo chunk
                current_chunk = para
                
                # Actualizar sección si parece título
                if len(para) < 100 and not para.endswith('.'):
                    current_section = para[:100]
        
        # Guardar último chunk
        if current_chunk:
            chunks.append({
                "text": current_chunk.strip(),
                "page": page_num,
                "section": current_section,
                "tokens": self.count_tokens(current_chunk)
            })
        
        return chunks
    
    def _extract_title(self, doc) -> str:
        """
        Intenta extraer el título del PDF
        """
        # 1. Intentar desde metadatos
        if doc.metadata.get("title"):
            title = doc.metadata["title"]
            if title and title.lower() not in ['untitled', 'sin título', 'documento']:
                return title
        
        # 2. Intentar desde primera página (primera línea de texto)
        try:
            first_page = doc[0].get_text("text")
            lines = [l.strip() for l in first_page.split('\n') if l.strip()]
            
            if lines:
                # Buscar la primera línea que parezca título (corta, sin punto final)
                for line in lines[:5]:  # Revisar primeras 5 líneas
                    if 10 < len(line) < 150 and not line.endswith('.'):
                        return line
                
                # Si no, usar primera línea
                return lines[0][:100]
        except:
            pass
        
        return "Documento sin título"
